import streamlit as st
import tempfile
import os
from typing import List, Dict
from langchain_community.document_loaders import PyPDFLoader, Docx2txtLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_google_genai import GoogleGenerativeAIEmbeddings, ChatGoogleGenerativeAI
from langchain_community.vectorstores import FAISS
from langchain.chains import ConversationalRetrievalChain
from langchain.schema import Document
import google.generativeai as genai
from pptx import Presentation

class CustomPPTLoader:
    def __init__(self, file_path: str, original_filename: str):
        self.file_path = file_path
        self.original_filename = original_filename

    def load(self) -> List[Document]:
        try:
            prs = Presentation(self.file_path)
            documents = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                text_content = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_content.append(shape.text.strip())
                
                if text_content:
                    text = "\n\n".join(text_content)
                    metadata = {"source": self.original_filename, "slide_number": slide_num}
                    documents.append(Document(page_content=text, metadata=metadata))
            
            return documents
        except Exception as e:
            st.error(f"Error processing PowerPoint file: {str(e)}")
            return []

class DocumentManager:
    def __init__(self):
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len
        )
        self.processed_files = []
        self.qa_chain = None
        self.chat_history = []

    def process_file(self, uploaded_file, progress_bar) -> List[Document]:
        try:
            original_filename = uploaded_file.name
            # Create a temporary file
            with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(original_filename)[1]) as tmp_file:
                tmp_file.write(uploaded_file.getvalue())
                file_path = tmp_file.name

            # Select appropriate loader and pass original filename
            if original_filename.endswith('.pdf'):
                loader = PyPDFLoader(file_path)
                documents = loader.load()
                # Update metadata with original filename
                for doc in documents:
                    doc.metadata['source'] = original_filename
            elif original_filename.endswith('.docx'):
                loader = Docx2txtLoader(file_path)
                documents = loader.load()
                # Update metadata with original filename
                for doc in documents:
                    doc.metadata['source'] = original_filename
            elif original_filename.endswith(('.pptx', '.ppt')):
                loader = CustomPPTLoader(file_path, original_filename)
                documents = loader.load()
            else:
                return []

            progress_bar.progress(0.5)
            
            if documents:
                chunks = self.text_splitter.split_documents(documents)
                self.processed_files.append(original_filename)
                progress_bar.progress(1.0)
                return chunks
            return []

        except Exception as e:
            st.error(f"Error processing {original_filename}: {str(e)}")
            return []
        finally:
            # Clean up temporary file
            if 'file_path' in locals():
                try:
                    os.unlink(file_path)
                except:
                    pass

    def setup_qa_system(self, uploaded_files):
        try:
            all_chunks = []
            progress_bars = {file.name: st.progress(0) for file in uploaded_files}
            status_text = st.empty()

            for uploaded_file in uploaded_files:
                status_text.text(f"Processing {uploaded_file.name}...")
                chunks = self.process_file(uploaded_file, progress_bars[uploaded_file.name])
                all_chunks.extend(chunks)

            if not all_chunks:
                st.error("No documents were successfully processed!")
                return False

            embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
            
            vectorstore = FAISS.from_documents(all_chunks, embeddings)
            
            self.qa_chain = ConversationalRetrievalChain.from_llm(
                ChatGoogleGenerativeAI(model="gemini-1.5-flash", temperature=0.7),
                vectorstore.as_retriever(),
                return_source_documents=True
            )
            return True

        except Exception as e:
            st.error(f"Error setting up QA system: {str(e)}")
            return False

    def ask_question(self, question: str) -> Dict:
        if not self.qa_chain:
            return {"error": "QA system not initialized. Please upload documents first."}
        try:
            result = self.qa_chain.invoke({
                "question": question,
                "chat_history": self.chat_history
            })
            
            # Use the original filenames from metadata
            sources = list({
                doc.metadata.get('source', 'Unknown source')
                for doc in result["source_documents"]
            })
            
            self.chat_history.append((question, result["answer"]))
            return {
                "answer": result["answer"],
                "sources": sources
            }
        except Exception as e:
            return {"error": f"Error processing question: {str(e)}"}

# Main function remains the same

def main():
    st.set_page_config(page_title="Document QA Assistant", page_icon="📚")
    st.title("Document QA Assistant 🤖")

    # Initialize session state
    if 'manager' not in st.session_state:
        st.session_state.manager = None
    if 'system_ready' not in st.session_state:
        st.session_state.system_ready = False
    if 'messages' not in st.session_state:
        st.session_state.messages = []

    # Get API key from environment variable or sidebar
    api_key = os.getenv('GOOGLE_API_KEY') or st.sidebar.text_input("Enter Google API Key", type="password")
    if api_key:
        os.environ["GOOGLE_API_KEY"] = api_key
        genai.configure(api_key=api_key)
    
    # File upload
    st.sidebar.header("Upload Documents")
    uploaded_files = st.sidebar.file_uploader(
        "Upload your documents",
        type=['pdf', 'docx', 'pptx', 'ppt'],
        accept_multiple_files=True
    )

    if uploaded_files and api_key and st.sidebar.button("Process Documents"):
        with st.spinner("Processing documents..."):
            manager = DocumentManager()
            if manager.setup_qa_system(uploaded_files):
                st.session_state.manager = manager
                st.session_state.system_ready = True
                st.sidebar.success("Ready to answer questions!")
            else:
                st.sidebar.error("Setup failed. Please try again.")

    # Chat interface
    if st.session_state.system_ready:
        if st.session_state.manager.processed_files:
            with st.expander("Processed Files"):
                for file in st.session_state.manager.processed_files:
                    st.text(f"✓ {file}")

        # Display chat history
        for message in st.session_state.messages:
            with st.chat_message(message["role"]):
                st.write(message["content"])
                if "sources" in message:
                    st.markdown("**Sources:**")
                    for source in message["sources"]:
                        st.markdown(f"- {source}")

        # Chat input
        if prompt := st.chat_input("Ask a question about your documents"):
            st.session_state.messages.append({"role": "user", "content": prompt})
            with st.chat_message("user"):
                st.write(prompt)

            with st.chat_message("assistant"):
                with st.spinner("Thinking..."):
                    response = st.session_state.manager.ask_question(prompt)
                    if "error" in response:
                        st.error(response["error"])
                    else:
                        st.write(response["answer"])
                        if response["sources"]:
                            st.markdown("**Sources:**")
                            for source in response["sources"]:
                                st.markdown(f"- {source}")
                        
                        st.session_state.messages.append({
                            "role": "assistant",
                            "content": response["answer"],
                            "sources": response["sources"]
                        })
    else:
        st.info("Upload your documents and provide API key to begin.")

if __name__ == "__main__":
    main()
