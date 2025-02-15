import os
import tempfile
from typing import List, Dict, Any
import streamlit as st
from langchain_community.document_loaders import PyPDFLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_google_genai import GoogleGenerativeAIEmbeddings, ChatGoogleGenerativeAI
from langchain_community.vectorstores import FAISS  # Using FAISS instead of Chroma
from langchain.chains import ConversationalRetrievalChain
import google.generativeai as genai
import docx2txt
from pptx import Presentation
from langchain.schema import Document

class CustomPPTLoader:
    """Custom PowerPoint loader that doesn't rely on NLTK"""
    def __init__(self, file_path):
        self.file_path = file_path

    def load(self) -> List[Document]:
        prs = Presentation(self.file_path)
        documents = []
        
        for i, slide in enumerate(prs.slides, 1):
            text_content = []
            
            # Extract text from shapes in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_content.append(shape.text)
            
            # Only create a document if there's actual content
            if text_content:
                text = "\n".join(text_content)
                metadata = {
                    "source": self.file_path,
                    "slide_number": i
                }
                documents.append(Document(
                    page_content=text,
                    metadata=metadata
                ))
        
        return documents

class DocumentManager:
    def __init__(self):
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len
        )
        self.processed_files = []
        self.qa_chain = None
        self.chat_history = []
        self.temp_dir = tempfile.mkdtemp()
        
    def process_file(self, file_path: str, progress_bar=None) -> List[str]:
        """Process a single file and return its chunks"""
        try:
            print(f"Processing: {file_path}")
            
            if file_path.endswith('.pdf'):
                loader = PyPDFLoader(file_path)
            elif file_path.endswith('.docx'):
                # Custom handling for Word documents
                text = docx2txt.process(file_path)
                documents = [Document(
                    page_content=text,
                    metadata={"source": file_path}
                )]
                return self.text_splitter.split_documents(documents)
            elif file_path.endswith(('.pptx', '.ppt')):
                loader = CustomPPTLoader(file_path)
            else:
                return []
                
            if progress_bar:
                progress_bar.progress(0.3)
                
            documents = loader.load()
            if progress_bar:
                progress_bar.progress(0.6)
            
            chunks = self.text_splitter.split_documents(documents)
            self.processed_files.append(os.path.basename(file_path))
            
            if progress_bar:
                progress_bar.progress(1.0)
            
            return chunks
        except Exception as e:
            st.error(f"Error processing {os.path.basename(file_path)}: {str(e)}")
            print(f"Detailed error: {str(e)}")  # For debugging
            return []

    def setup_qa_system(self, uploaded_files):
        """Initialize the QA system with uploaded documents"""
        try:
            all_chunks = []
            
            # Create progress containers for each file
            progress_bars = {file.name: st.progress(0) for file in uploaded_files}
            status_text = st.empty()
            
            for uploaded_file in uploaded_files:
                status_text.text(f"Processing {uploaded_file.name}...")
                
                # Save uploaded file to temporary directory
                temp_path = os.path.join(self.temp_dir, uploaded_file.name)
                with open(temp_path, "wb") as f:
                    f.write(uploaded_file.getvalue())
                
                # Process the file
                chunks = self.process_file(temp_path, progress_bars[uploaded_file.name])
                all_chunks.extend(chunks)
                
                # Clean up temporary file
                os.remove(temp_path)
            
            if not all_chunks:
                st.error("No documents were successfully processed!")
                return False
            
            status_text.text("Initializing embeddings...")
            embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
            
            status_text.text("Creating vector store...")
            vectorstore = FAISS.from_documents(
                documents=all_chunks,
                embedding=embeddings
            )
            
            status_text.text("Setting up QA chain...")
            self.qa_chain = ConversationalRetrievalChain.from_llm(
                ChatGoogleGenerativeAI(model="gemini-1.5-flash", temperature=0.7),
                vectorstore.as_retriever(),
                return_source_documents=True
            )
            
            status_text.text("Setup complete!")
            return True
            
        except Exception as e:
            st.error(f"Error setting up QA system: {str(e)}")
            print(f"Detailed error: {str(e)}")  # For debugging
            return False

    def ask_question(self, question: str) -> Dict:
        """Ask a question and get a response with source information"""
        if not self.qa_chain:
            return {"error": "QA system not initialized. Please upload documents first."}
        try:
            result = self.qa_chain.invoke({
                "question": question,
                "chat_history": self.chat_history
            })
            
            sources = list({
                doc.metadata.get('source', 'Unknown source')
                for doc in result["source_documents"]
            })
            
            self.chat_history.append((question, result["answer"]))
            return {
                "answer": result["answer"],
                "sources": sources
            }
        except Exception as e:
            return {"error": f"Error processing question: {str(e)}"}

# Rest of your main() function remains the same...

def main():
    st.set_page_config(page_title="FolderFlow Document Assistant", page_icon="📚")
    st.title("FolderFlow Document Assistant 🤖")

    # Initialize session state variables
    if 'manager' not in st.session_state:
        st.session_state.manager = None
    if 'system_ready' not in st.session_state:
        st.session_state.system_ready = False
    if 'messages' not in st.session_state:
        st.session_state.messages = []

    # Configure Google API
    api_key = st.sidebar.text_input("Enter Google API Key", type="password")
    if api_key:
        os.environ["GOOGLE_API_KEY"] = api_key
        genai.configure(api_key=api_key)
    else:
        st.info("Please enter your Google API key in the sidebar to continue.")
        return

    # File upload section
    st.sidebar.header("Upload Documents")
    uploaded_files = st.sidebar.file_uploader(
        "Drop your documents here or click to upload",
        type=['pdf', 'docx', 'pptx', 'ppt'],
        accept_multiple_files=True,
        help="Supported formats: PDF, DOCX, PPTX, PPT"
    )

    # Initialize system button
    if uploaded_files and st.sidebar.button("Process Documents"):
        with st.spinner("Setting up the document management system..."):
            manager = DocumentManager()
            if manager.setup_qa_system(uploaded_files):
                st.session_state.manager = manager
                st.session_state.system_ready = True
                st.sidebar.success("System initialized successfully!")
            else:
                st.sidebar.error("Failed to initialize the system. Please check your files.")

    # Main chat interface
    if st.session_state.system_ready:
        # Display processed files
        if st.session_state.manager.processed_files:
            with st.expander("Processed Files"):
                for file in st.session_state.manager.processed_files:
                    st.text(f"✓ {file}")

        # Chat interface
        st.markdown("### Ask me anything about your documents!")
        
        # Display chat messages
        for message in st.session_state.messages:
            with st.chat_message(message["role"]):
                st.write(message["content"])
                if "sources" in message:
                    st.markdown("**Sources:**")
                    unique_sources = {os.path.basename(source) for source in message["sources"]}
                    for source in unique_sources:
                        st.markdown(f"- {source}")

        # Chat input
        if prompt := st.chat_input("Your question"):
            # Add user message to chat history
            st.session_state.messages.append({"role": "user", "content": prompt})
            
            # Display user message
            with st.chat_message("user"):
                st.write(prompt)

            # Generate and display assistant response
            with st.chat_message("assistant"):
                with st.spinner("Thinking..."):
                    response = st.session_state.manager.ask_question(prompt)
                    if "error" in response:
                        st.error(response["error"])
                    else:
                        st.write(response["answer"])
                        if response["sources"]:
                            st.markdown("**Sources:**")
                            unique_sources = {os.path.basename(source) for source in response["sources"]}
                            for source in unique_sources:
                                st.markdown(f"- {source}")
                        
                        # Add assistant response to chat history
                        st.session_state.messages.append({
                            "role": "assistant",
                            "content": response["answer"],
                            "sources": response["sources"]
                        })
    else:
        st.info("Please upload your documents and initialize the system using the sidebar controls.")

if __name__ == "__main__":
    main()
